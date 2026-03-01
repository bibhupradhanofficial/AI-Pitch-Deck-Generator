from __future__ import annotations

import asyncio
import os
import tempfile
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

from models.schemas import SlideSpec


@dataclass(frozen=True)
class PptxServiceConfig:
    title_font_size_pt: int = 36
    body_font_size_pt: int = 20


class PptxService:
    def __init__(self, config: PptxServiceConfig | None = None) -> None:
        self._config = config or PptxServiceConfig()

    def build_deck(self, slides: list[SlideSpec], out_path: Path, chart_paths: dict[int, Path] | None = None) -> Path:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()

        for idx, spec in enumerate(slides):
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            title = slide.shapes.title
            body = slide.placeholders[1]

            title.text = spec.title
            title.text_frame.paragraphs[0].font.size = Pt(self._config.title_font_size_pt)

            tf = body.text_frame
            tf.clear()
            if spec.bullets:
                p0 = tf.paragraphs[0]
                p0.text = spec.bullets[0]
                p0.font.size = Pt(self._config.body_font_size_pt)
                for bullet in spec.bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    p.font.size = Pt(self._config.body_font_size_pt)
            else:
                tf.text = ""

            if spec.notes:
                slide.notes_slide.notes_text_frame.text = spec.notes

            if chart_paths and idx in chart_paths:
                img_path = str(chart_paths[idx])
                left = Inches(6.4)
                top = Inches(1.6)
                height = Inches(3.2)
                slide.shapes.add_picture(img_path, left, top, height=height)

        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(out_path)
        return out_path


def _get_tmp_dir() -> Path:
    tmp_dir = Path("/tmp")
    if tmp_dir.exists():
        return tmp_dir
    return Path(tempfile.gettempdir())


def _coerce_bullets(body: object) -> list[str]:
    if body is None:
        return []
    if isinstance(body, list):
        return [str(x).strip() for x in body if str(x).strip()]
    if isinstance(body, str):
        return [line.strip() for line in body.splitlines() if line.strip()]
    return [str(body).strip()] if str(body).strip() else []


def _add_themed_slide(prs, *, title: str, cover: bool) -> object:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    dark_bg = RGBColor(0x0F, 0x0F, 0x1A)
    accent = RGBColor(0x66, 0x7E, 0xEA)
    white = RGBColor(0xFF, 0xFF, 0xFF)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = dark_bg
    bg.line.fill.background()

    accent_bar_h = Inches(4 / 96)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, accent_bar_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    if cover:
        title_left = 0
        title_top = Inches(2.6)
        title_w = prs.slide_width
        title_h = Inches(1.4)
        title_pt = 48
        title_align = PP_ALIGN.CENTER
    else:
        title_left = Inches(0.8)
        title_top = Inches(0.6)
        title_w = prs.slide_width - Inches(1.6)
        title_h = Inches(0.9)
        title_pt = 36
        title_align = PP_ALIGN.LEFT

    title_box = slide.shapes.add_textbox(title_left, title_top, title_w, title_h)
    tf = title_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = title_align
    p.font.size = Pt(title_pt)
    p.font.bold = True
    p.font.color.rgb = white

    return slide


def _add_body_bullets(prs, slide, bullets: list[str], *, cover: bool) -> None:
    if not bullets:
        return

    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    white = RGBColor(0xFF, 0xFF, 0xFF)

    if cover:
        left = Inches(1.5)
        top = Inches(4.0)
        width = prs.slide_width - Inches(3.0)
        height = prs.slide_height - top - Inches(0.8)
        align = PP_ALIGN.CENTER
    else:
        left = Inches(1.0)
        top = Inches(1.6)
        width = prs.slide_width - Inches(2.0)
        height = prs.slide_height - top - Inches(0.8)
        align = PP_ALIGN.LEFT

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸ {bullet}"
        p.alignment = align
        p.font.size = Pt(18)
        p.font.color.rgb = white


def _add_speaker_notes(slide, notes: object) -> None:
    text = "" if notes is None else str(notes)
    slide.notes_slide.notes_text_frame.text = text


def _add_image_fitted(slide, img_bytes: bytes, *, left: int, top: int, box_w: int, box_h: int) -> None:
    from PIL import Image

    with Image.open(BytesIO(img_bytes)) as im:
        px_w, px_h = im.size
    if px_w <= 0 or px_h <= 0:
        return

    scale = min(box_w / px_w, box_h / px_h)
    target_w = max(1, int(px_w * scale))
    target_h = max(1, int(px_h * scale))
    pic_left = left + (box_w - target_w) // 2
    pic_top = top + (box_h - target_h) // 2

    slide.shapes.add_picture(BytesIO(img_bytes), pic_left, pic_top, width=target_w, height=target_h)


def _add_charts_slides(prs, chart_images: list[bytes]) -> None:
    if not chart_images:
        return

    from pptx.util import Inches

    margin_x = Inches(0.7)
    gap_x = Inches(0.4)
    top = Inches(1.5)
    bottom = Inches(0.6)
    box_h = prs.slide_height - top - bottom
    box_w = (prs.slide_width - (2 * margin_x) - gap_x) // 2

    for i in range(0, len(chart_images), 2):
        slide = _add_themed_slide(prs, title="Charts", cover=False)
        pair = chart_images[i : i + 2]
        for j, img in enumerate(pair):
            left = margin_x + (box_w + gap_x) * j
            _add_image_fitted(slide, img, left=left, top=top, box_w=box_w, box_h=box_h)


def _add_product_visuals_slides(prs, images: list[bytes]) -> None:
    if not images:
        return

    from pptx.util import Inches

    margin = Inches(0.7)
    gap = Inches(0.3)
    top = Inches(1.5)
    bottom = Inches(0.6)

    max_per_slide = 4
    grid_w = prs.slide_width - 2 * margin
    grid_h = prs.slide_height - top - bottom

    for i in range(0, len(images), max_per_slide):
        slide_images = images[i : i + max_per_slide]
        slide = _add_themed_slide(prs, title="Product Visuals", cover=False)

        if len(slide_images) == 1:
            _add_image_fitted(slide, slide_images[0], left=margin, top=top, box_w=grid_w, box_h=grid_h)
            continue

        if len(slide_images) == 2:
            cell_w = (grid_w - gap) // 2
            for j, img in enumerate(slide_images):
                left = margin + (cell_w + gap) * j
                _add_image_fitted(slide, img, left=left, top=top, box_w=cell_w, box_h=grid_h)
            continue

        cell_w = (grid_w - gap) // 2
        cell_h = (grid_h - gap) // 2
        positions = [
            (margin, top),
            (margin + cell_w + gap, top),
            (margin, top + cell_h + gap),
            (margin + cell_w + gap, top + cell_h + gap),
        ]
        for (left, cell_top), img in zip(positions, slide_images, strict=False):
            _add_image_fitted(slide, img, left=left, top=cell_top, box_w=cell_w, box_h=cell_h)


def _render_pptx(slides: list[dict], chart_images: list[bytes], product_images: list[bytes], out_path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for idx, s in enumerate(slides):
        title = str(s.get("title") or "").strip()
        body = _coerce_bullets(s.get("body"))
        notes = s.get("speaker_notes") if "speaker_notes" in s else s.get("notes")

        slide = _add_themed_slide(prs, title=title, cover=(idx == 0))
        _add_body_bullets(prs, slide, body, cover=(idx == 0))
        _add_speaker_notes(slide, notes)

    _add_charts_slides(prs, chart_images)
    _add_product_visuals_slides(prs, product_images)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


async def build_pptx(slides: list[dict], chart_urls: list[str], image_urls: list[str], session_id: str) -> str:
    import requests

    async def _download(url: str) -> bytes:
        def _get() -> bytes:
            resp = requests.get(url, timeout=60)
            resp.raise_for_status()
            return resp.content

        return await asyncio.to_thread(_get)

    chart_images, product_images = await asyncio.gather(
        asyncio.gather(*[_download(u) for u in (chart_urls or [])]),
        asyncio.gather(*[_download(u) for u in (image_urls or [])]),
    )

    tmp_dir = _get_tmp_dir()
    out_path = tmp_dir / f"{session_id}_pitch_deck.pptx"
    await asyncio.to_thread(_render_pptx, slides, list(chart_images), list(product_images), out_path)

    bucket = (os.getenv("GCS_BUCKET_NAME") or "").strip()
    if not bucket:
        raise RuntimeError("GCS_BUCKET_NAME is not configured")
    prefix = (os.getenv("GCS_PREFIX") or "pitch-decks").strip().strip("/")
    destination_blob_name = f"{prefix}/{out_path.name}"

    from services.gcs_service import upload_file

    url = await asyncio.to_thread(upload_file, str(out_path), destination_blob_name)
    return url
